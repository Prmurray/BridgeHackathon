import json
import os
from parserr import parse_files  # Import parse_files from parserr.py
from pptx import Presentation  # Import the Presentation class from python-pptx

def save_parsed_data(parsed_data, filename='parsed_data.json'):
    """
    Save parsed data to a JSON file.
    """
    with open(filename, 'w', encoding='utf-8') as f:
        json.dump(parsed_data, f, indent=4, ensure_ascii=False)
    print(f"Parsed data saved to {filename}.")

def load_parsed_data(filename='parsed_data.json'):
    """
    Load parsed data from a JSON file and handle any potential errors.
    """
    if not os.path.exists(filename):
        print(f"File {filename} does not exist.")
        return []

    with open(filename, 'r', encoding='utf-8') as f:
        try:
            parsed_data = json.load(f)
            if not parsed_data:
                print(f"File {filename} is empty or contains no data.")
                return []
        except json.JSONDecodeError as e:
            print(f"Error loading JSON from {filename}: {e}")
            return []

    return parsed_data

def get_pptx_files(directory='profiles'):
    """
    Retrieve all .pptx files from the specified directory.
    """
    if not os.path.exists(directory):
        print(f"Directory {directory} does not exist.")
        return []

    return [os.path.join(directory, file) for file in os.listdir(directory) if file.endswith('.pptx')]

def extract_slides_from_pptx(pptx_file):
    """
    Extract text content from each slide of a PowerPoint file.
    """
    prs = Presentation(pptx_file)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_data.append({
            'slide_num': slide_num,
            'raw_data': ' '.join(slide_text)  # Combine all text parts from the slide
        })

    return slides_data

def read_pptx_files(pptx_files):
    """
    Read the contents of each .pptx file and extract text from slides.
    """
    files_data = []
    for pptx_file in pptx_files:
        slides_data = extract_slides_from_pptx(pptx_file)
        files_data.append({
            'filename': os.path.basename(pptx_file),
            'slides': slides_data
        })
    return files_data

def main():
    # Get all .pptx files from the 'profiles' directory
    pptx_files = get_pptx_files('profiles')

    if not pptx_files:
        print("No PowerPoint files found in the 'profiles' directory.")
        return

    # Read content from PowerPoint files
    files_data = read_pptx_files(pptx_files)

    # Parse each PowerPoint file and add the parsed data to the list
    all_parsed_data = parse_files(files_data)

    # Save the parsed data to parsed_data.json
    save_parsed_data(all_parsed_data)

if __name__ == "__main__":
    main()
