import os
import json
from pptx import Presentation
from parserr import parse_presentation, clean_text
from db_manager import connect_db, create_tables, store_consultant_info, store_slide_data

def extract_text_from_slide(slide) -> str:
    """
    Extracts text from a single slide in a .pptx file.
    """
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text.append(shape.text)
    return clean_text(" ".join(slide_text))

def parse_pptx_file(filepath: str):
    """
    Parses a .pptx file and extracts text from each slide.
    """
    prs = Presentation(filepath)
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = extract_text_from_slide(slide)
        slides_data.append({'slide_num': slide_num, 'raw_data': slide_text})

    return parse_presentation(slides_data)

def parse_ppt_files(directory: str):
    """
    Parses all .pptx files in the specified directory and extracts relevant data.
    """
    parsed_files_data = []

    for filename in os.listdir(directory):
        if filename.endswith(".pptx"):
            pptx_path = os.path.join(directory, filename)
            parsed_file_data = {
                "filename": filename,
                "slides": parse_pptx_file(pptx_path)
            }
            parsed_files_data.append(parsed_file_data)

    return parsed_files_data

def save_parsed_data_to_json(parsed_data, output_file='parsed_data.json'):
    """
    Saves the parsed data to a JSON file.
    """
    with open(output_file, 'w', encoding='utf-8') as f:
        json.dump(parsed_data, f, ensure_ascii=False, indent=4)

def insert_parsed_data_to_db(parsed_data):
    """
    Insert parsed data into the Azure SQL Database.
    """
    conn = connect_db()
    if conn:
        create_tables(conn)  # Ensure tables exist before inserting data

        for file_data in parsed_data:
            for slide in file_data['slides']:
                parsed_data = slide['parsed_data']
                name = parsed_data.get('name', '')
                title = parsed_data.get('title', '')
                email = parsed_data.get('email', '')
                mobile = parsed_data.get('mobile', '')
                location = parsed_data.get('location', '')
                data = parsed_data.get('data', '')

                consultant_info = {
                    'name': name,
                    'title': title,
                    'mobile': mobile,
                    'location': location,
                    'email': email
                }

                # Store consultant information and retrieve the consultant ID
                consultant_id = store_consultant_info(conn, consultant_info)

                # Store slide data linked to the consultant ID
                slide_data = {
                    'slide_id': slide['slide_num'],  # Assuming slide number as slide_id
                    'data': data
                }
                store_slide_data(conn, consultant_id, [slide_data])

        conn.close()
        print("Data successfully inserted into Azure SQL Database.")
    else:
        print("Failed to connect to the database.")

def main():
    """
    Main function to execute parsing, saving to JSON, and sending to Azure SQL Database.
    """
    directory = 'profiles'  # Update this to your directory path
    parsed_data = parse_ppt_files(directory)

    # Save parsed data to JSON (optional)
    save_parsed_data_to_json(parsed_data)

    # Insert parsed data to Azure SQL Database
    insert_parsed_data_to_db(parsed_data)

if __name__ == "__main__":
    main()
