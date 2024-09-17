from pptx import Presentation
from pathlib import Path
import re
from sqlalchemy import create_engine, Table, Column, Integer, String, MetaData
from sqlalchemy.orm import sessionmaker
from config import get_db_connection_string  # Use the config file for DB connection
import fitz  # PyMuPDF for PDF handling

# Set up Azure SQL connection using SQLAlchemy
DATABASE_URL = get_db_connection_string()
engine = create_engine(DATABASE_URL)
Session = sessionmaker(bind=engine)
session = Session()

metadata = MetaData()

# Define the consultants table
consultants_table = Table(
    'consultants', metadata,
    Column('id', Integer, primary_key=True),
    Column('name', String(255)),
    Column('profile_text', String)
)

metadata.create_all(engine)

# Directory containing PowerPoint and PDF profiles
directory = Path('./profiles/')

def clean_text(text: str) -> str:
    """
    Clean extracted text by removing unwanted characters, fixing common OCR errors, and trimming spaces.
    """
    text = re.sub(r'[\u000b\n]+', ' ', text)  # Replace unwanted characters with space
    text = re.sub(r'\s+', ' ', text).strip()  # Remove extra spaces and trim
    return text

# Function to extract consultant name from file name or fallback to text extraction
def extract_consultant_name(file_name, text):
    """
    Extract consultant name using file name patterns or fallback to generic pattern in text.
    """
    # First attempt: Extract consultant name from file name pattern
    match = re.match(r"(\w+)\s+(\w+)\s+-\s+Consultant Profile", file_name)
    if match:
        first_name, last_name = match.groups()
        return f"{first_name} {last_name}"
    
    # Fallback 1: Handle names with mixed cases (e.g., McDonald, DiPrima)
    complex_name_pattern = r'\b[A-Z][a-z]+(?:[a-z]*[A-Z][a-z]*)?\s[A-Z][a-z]+(?:[a-z]*[A-Z][a-z]*)?\b'
    complex_name_match = re.search(complex_name_pattern, text)
    if complex_name_match:
        return complex_name_match.group(0)

    # Fallback 2: General pattern for names in the text (First Last)
    name_pattern = r'\b[A-Z][a-z]+\s[A-Z][a-z]+\b'
    name_match = re.search(name_pattern, text)
    if name_match:
        return name_match.group(0)

    # If no match is found, return "Unknown"
    return "Unknown"
    
# Function to extract text from PowerPoint presentation
def extract_text_from_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text_runs = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    cleaned = clean_text(shape.text.strip())
                    text_runs.append(cleaned)
        
        return "\n".join(text_runs)
    except Exception as e:
        print(f"Error processing PowerPoint file {pptx_path}: {e}")
        return ""

# Function to extract text from PDF file
def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)  # Open the PDF file
        text_runs = []

        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            text_runs.append(page.get_text())

        doc.close()
        return clean_text(" ".join(text_runs))
    except Exception as e:
        print(f"Error processing PDF file {pdf_path}: {e}")
        return ""

def load_profiles_to_database():
    files = directory.iterdir()

    for file in files:
        # Handle PowerPoint files (PPTX, PPTM)
        if file.suffix.lower() in ['.pptx', '.pptm']:
            # Extract text from the file
            raw_text = extract_text_from_pptx(file)

        # Handle PDF files
        elif file.suffix.lower() == '.pdf':
            # Extract text from the file
            raw_text = extract_text_from_pdf(file)

        # Skip other file types
        else:
            print(f"Skipping unsupported file type: {file.name}")
            continue
        
        if not raw_text:  # Skip if the text extraction failed
            continue

        # Extract the consultant name using the combined approach
        consultant_name = extract_consultant_name(file.name, raw_text)

        # Clean the extracted text
        cleaned_text = raw_text.replace("\n", " ").replace("\t", " ").strip()  # Basic cleaning

        # Insert profile data into the database
        insert_query = consultants_table.insert().values(name=consultant_name, profile_text=cleaned_text)
        session.execute(insert_query)
    
    session.commit()

if __name__ == "__main__":
    load_profiles_to_database()
    print("Profiles have been loaded into the database.")
